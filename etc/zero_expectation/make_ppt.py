"""
Zero-Expectation QUBO 방법론 PPT 생성 스크립트
- 연구실 세미나용, PDF 논리 흐름 기반
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── 색상 팔레트 ──
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x1A, 0x1A, 0x2E)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
MID_GRAY = RGBColor(0x66, 0x66, 0x66)
LIGHT_GRAY = RGBColor(0xF0, 0xF0, 0xF5)
ACCENT_BLUE = RGBColor(0x2D, 0x5B, 0xE3)
ACCENT_RED = RGBColor(0xE0, 0x3E, 0x3E)
ACCENT_GREEN = RGBColor(0x27, 0xAE, 0x60)
ACCENT_ORANGE = RGBColor(0xF3, 0x9C, 0x12)
SOFT_BLUE_BG = RGBColor(0xEB, 0xF0, 0xFF)
SOFT_RED_BG = RGBColor(0xFF, 0xEB, 0xEB)
SOFT_GREEN_BG = RGBColor(0xE8, 0xF8, 0xEE)
SOFT_YELLOW_BG = RGBColor(0xFF, 0xF8, 0xE1)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# ── 헬퍼 함수 ──

def add_rounded_rect(slide, left, top, width, height, fill_color, border_color=None, border_width=Pt(0)):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()
    # 둥근 정도 조절
    shape.adjustments[0] = 0.05
    return shape


def add_rect(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height):
    return slide.shapes.add_textbox(left, top, width, height)


def set_text(tf, text, font_size=14, color=DARK_GRAY, bold=False, italic=False,
             alignment=PP_ALIGN.LEFT, font_name="맑은 고딕", spacing_after=Pt(4)):
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.italic = italic
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = spacing_after
    return p


def add_paragraph(tf, text, font_size=14, color=DARK_GRAY, bold=False, italic=False,
                  alignment=PP_ALIGN.LEFT, font_name="맑은 고딕", spacing_after=Pt(4),
                  spacing_before=Pt(0), level=0):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.italic = italic
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = spacing_after
    p.space_before = spacing_before
    p.level = level
    return p


def add_rich_paragraph(tf, runs, alignment=PP_ALIGN.LEFT, spacing_after=Pt(4),
                       spacing_before=Pt(0)):
    """runs = [(text, font_size, color, bold, italic, font_name), ...]"""
    p = tf.add_paragraph()
    p.alignment = alignment
    p.space_after = spacing_after
    p.space_before = spacing_before
    for run_data in runs:
        text, fs, clr, bld, itl, fn = run_data
        run = p.add_run()
        run.text = text
        run.font.size = Pt(fs)
        run.font.color.rgb = clr
        run.font.bold = bld
        run.font.italic = itl
        run.font.name = fn
    return p


def add_slide_number(slide, num, total):
    tb = add_text_box(slide, Inches(12.3), Inches(7.05), Inches(0.8), Inches(0.35))
    set_text(tb.text_frame, f"{num} / {total}", font_size=10, color=MID_GRAY,
             alignment=PP_ALIGN.RIGHT)


def add_section_header(slide, left, top, width, text, color=ACCENT_BLUE):
    """작은 섹션 라벨"""
    tb = add_text_box(slide, left, top, width, Inches(0.35))
    set_text(tb.text_frame, text, font_size=11, color=color, bold=True,
             font_name="맑은 고딕")
    # 밑줄 라인
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top + Inches(0.28),
                                   Inches(0.6), Pt(2.5))
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()


# ────────────────────────────────────────────────────────────────
# 슬라이드 0: 표지
# ────────────────────────────────────────────────────────────────

slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

# 배경
bg = slide.background.fill
bg.solid()
bg.fore_color.rgb = WHITE

# 상단 파란 바
add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), ACCENT_BLUE)

# 왼쪽 장식 바
add_rect(slide, Inches(0.8), Inches(2.0), Pt(5), Inches(2.8), ACCENT_BLUE)

# 제목
tb = add_text_box(slide, Inches(1.2), Inches(2.2), Inches(10), Inches(1.0))
set_text(tb.text_frame, "Zero-Expectation QUBO", font_size=40, color=BLACK,
         bold=True, font_name="맑은 고딕")
add_paragraph(tb.text_frame, "정답을 심은 QUBO를 랜덤과 구별 불가능하게 만들기",
              font_size=22, color=MID_GRAY, bold=False)

# 핵심 질문 박스
box = add_rounded_rect(slide, Inches(1.2), Inches(4.0), Inches(10.5), Inches(1.0),
                        SOFT_BLUE_BG, ACCENT_BLUE, Pt(1.5))
tb = add_text_box(slide, Inches(1.5), Inches(4.1), Inches(10), Inches(0.8))
tf = tb.text_frame
tf.word_wrap = True
set_text(tf, "핵심 질문", font_size=13, color=ACCENT_BLUE, bold=True)
add_paragraph(tf, "임의의 n-bit x*가 optimal인 QUBO를 만들 수 있을까?  단, x*를 비밀로 감췄을 때, 이를 찾는 게 어려워야 함.",
              font_size=16, color=DARK_GRAY)

# 하단 정보
tb = add_text_box(slide, Inches(1.2), Inches(5.8), Inches(10), Inches(0.5))
set_text(tb.text_frame, "연구실 세미나", font_size=14, color=MID_GRAY)

add_slide_number(slide, 0, 3)


# ────────────────────────────────────────────────────────────────
# 슬라이드 1: 동기 — 정답을 아는 QUBO가 필요하다
# ────────────────────────────────────────────────────────────────

slide = prs.slides.add_slide(prs.slide_layouts[6])
bg = slide.background.fill
bg.solid()
bg.fore_color.rgb = WHITE

add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), ACCENT_BLUE)

# 슬라이드 제목
tb = add_text_box(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.6))
set_text(tb.text_frame, "1. 동기 — 정답을 아는 QUBO가 필요하다", font_size=26,
         color=BLACK, bold=True)

# ── 좌측: QUBO 정의 + 솔버 평가 문제 ──

# QUBO 정의 박스
add_section_header(slide, Inches(0.6), Inches(1.1), Inches(3), "QUBO 정의")

box = add_rounded_rect(slide, Inches(0.6), Inches(1.5), Inches(5.8), Inches(1.6),
                        LIGHT_GRAY)
tb = add_text_box(slide, Inches(0.9), Inches(1.6), Inches(5.3), Inches(1.4))
tf = tb.text_frame
tf.word_wrap = True
set_text(tf, "n×n 실수 행렬 Q가 주어졌을 때,", font_size=14, color=DARK_GRAY)
add_paragraph(tf, "", font_size=6, color=DARK_GRAY)
add_paragraph(tf, "    min  Σᵢ Σⱼ qᵢ,ⱼ · xᵢ · xⱼ      (xᵢ ∈ {0, 1})",
              font_size=16, color=BLACK, bold=True, font_name="Consolas")
add_paragraph(tf, "", font_size=6, color=DARK_GRAY)
add_paragraph(tf, "→ classic 환경에서 optimal을 찾는 효율적 알고리즘 없음 (NP-hard)",
              font_size=13, color=ACCENT_RED, bold=True)

# 솔버 평가 문제
add_section_header(slide, Inches(0.6), Inches(3.4), Inches(5), "솔버 평가 문제")

tb = add_text_box(slide, Inches(0.6), Inches(3.8), Inches(5.8), Inches(1.8))
tf = tb.text_frame
tf.word_wrap = True
set_text(tf, "누군가 QUBO 솔버 A를 제안 → 어떻게 평가할까?", font_size=15,
         color=DARK_GRAY, bold=True)
add_paragraph(tf, "", font_size=4)
add_paragraph(tf, "• optimal을 확정지을 수 없음 (그 방법이 있으면 솔버가 필요 없음)",
              font_size=13, color=DARK_GRAY)
add_paragraph(tf, "• 다만 A를 평가할 데이터셋은 만들 수 있을지도 모름", font_size=13,
              color=DARK_GRAY)
add_paragraph(tf, "", font_size=6)
add_paragraph(tf, "비유: 다항식의 해를 찾는 방법의 정확성 측정", font_size=14,
              color=ACCENT_BLUE, bold=True)
add_paragraph(tf, "  {x² + 4x + 4 = 0,  해: {+2, −2}} 같은 데이터가 필요",
              font_size=13, color=DARK_GRAY, font_name="Consolas")
add_paragraph(tf, "  50개의 해를 아는 50차 다항식: a·∏(x − sᵢ) = 0 으로 생성 가능",
              font_size=13, color=DARK_GRAY, font_name="Consolas")

# ── 우측: 단순 구성의 한계 + 핵심 조건 ──

add_section_header(slide, Inches(7.0), Inches(1.1), Inches(5), "단순 구성의 한계")

# 쉬운 예시 박스
box = add_rounded_rect(slide, Inches(7.0), Inches(1.5), Inches(5.8), Inches(1.8),
                        SOFT_RED_BG, ACCENT_RED, Pt(1))
tb = add_text_box(slide, Inches(7.3), Inches(1.6), Inches(5.3), Inches(1.6))
tf = tb.text_frame
tf.word_wrap = True
set_text(tf, "optimal = 1011인 QUBO를 만들고 싶다면?", font_size=14,
         color=DARK_GRAY, bold=True)
add_paragraph(tf, "", font_size=4)
add_paragraph(tf, "    y  =  −x₀ + x₁ − x₂ − x₃", font_size=18, color=ACCENT_RED,
              bold=True, font_name="Consolas")
add_paragraph(tf, "", font_size=4)
add_paragraph(tf, "→ 2차 항이 없는 trivial한 QUBO", font_size=13, color=DARK_GRAY)
add_paragraph(tf, "→ bᵢ = 0이면 계수 양수, bᵢ = 1이면 계수 음수 ← 정답이 노출됨!",
              font_size=13, color=ACCENT_RED, bold=True)

# 핵심 조건 박스
add_section_header(slide, Inches(7.0), Inches(3.6), Inches(5), "핵심 조건",
                   color=ACCENT_GREEN)

box = add_rounded_rect(slide, Inches(7.0), Inches(4.0), Inches(5.8), Inches(1.5),
                        SOFT_GREEN_BG, ACCENT_GREEN, Pt(1.5))
tb = add_text_box(slide, Inches(7.3), Inches(4.1), Inches(5.3), Inches(1.3))
tf = tb.text_frame
tf.word_wrap = True
set_text(tf, "만들어진 문제가 랜덤 QUBO와 비슷한 난이도여야 함",
         font_size=15, color=DARK_GRAY, bold=True)
add_paragraph(tf, "(구별 불가해야 함)", font_size=15, color=ACCENT_GREEN, bold=True)
add_paragraph(tf, "", font_size=6)
add_paragraph(tf, "• Q 행렬의 계수 분포를 보고 target을 추론할 수 없어야 함",
              font_size=13, color=DARK_GRAY)
add_paragraph(tf, "• 즉, Q의 각 계수의 기댓값이 target에 의존하지 않아야 함",
              font_size=13, color=DARK_GRAY)
add_paragraph(tf, "• 목표: E[qᵢⱼ] = 0,  E[qᵢᵢ] = target bit와 무관",
              font_size=14, color=ACCENT_GREEN, bold=True, font_name="Consolas")

# 하단 요약 화살표
box = add_rounded_rect(slide, Inches(0.6), Inches(5.9), Inches(12.2), Inches(0.7),
                        SOFT_BLUE_BG, ACCENT_BLUE, Pt(1))
tb = add_text_box(slide, Inches(0.9), Inches(5.95), Inches(11.8), Inches(0.6))
tf = tb.text_frame
set_text(tf,
    "정리:  정답 x*를 심은 QUBO를 만들되, Q 행렬만 봐서는 x*에 대한 정보를 얻을 수 없도록  →  E[계수] = 0이 핵심",
    font_size=14, color=ACCENT_BLUE, bold=True, alignment=PP_ALIGN.CENTER)

add_slide_number(slide, 1, 3)


# ────────────────────────────────────────────────────────────────
# 슬라이드 2: 구성법 — Step 1 & Step 2
# ────────────────────────────────────────────────────────────────

slide = prs.slides.add_slide(prs.slide_layouts[6])
bg = slide.background.fill
bg.solid()
bg.fore_color.rgb = WHITE

add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), ACCENT_BLUE)

tb = add_text_box(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.6))
set_text(tb.text_frame,
         "2. 구성법 — x* = b₀b₁···bₙ₋₁ 이 해인 QUBO 만들기",
         font_size=26, color=BLACK, bold=True)

# ── Step 1 ──

add_section_header(slide, Inches(0.6), Inches(1.1), Inches(3), "Step 1. 1차 항 (대각)")

box = add_rounded_rect(slide, Inches(0.6), Inches(1.5), Inches(5.8), Inches(2.3),
                        SOFT_BLUE_BG, ACCENT_BLUE, Pt(1))
tb = add_text_box(slide, Inches(0.9), Inches(1.6), Inches(5.3), Inches(2.1))
tf = tb.text_frame
tf.word_wrap = True
set_text(tf, "임의의 양수 a₀, a₁, ..., aₙ₋₁ 을 뽑고,", font_size=14, color=DARK_GRAY)
add_paragraph(tf, "", font_size=6)
add_paragraph(tf, "  bᵢ = 0 이면  → +aᵢxᵢ  를 목적함수에 더함", font_size=15,
              color=DARK_GRAY, font_name="Consolas")
add_paragraph(tf, "  bᵢ = 1 이면  → −aᵢxᵢ  를 목적함수에 더함", font_size=15,
              color=DARK_GRAY, font_name="Consolas")
add_paragraph(tf, "", font_size=8)
add_paragraph(tf, "원리:", font_size=13, color=ACCENT_BLUE, bold=True)
add_paragraph(tf, "  bᵢ=0 → xᵢ=0이 이득 (+aᵢ·0 < +aᵢ·1)", font_size=13,
              color=DARK_GRAY)
add_paragraph(tf, "  bᵢ=1 → xᵢ=1이 이득 (−aᵢ·1 < −aᵢ·0)", font_size=13,
              color=DARK_GRAY)
add_paragraph(tf, "", font_size=4)
add_paragraph(tf, "⚠ 단, 이것만으로는 1차 항뿐이라 쉽게 풀림 → Step 2 필요",
              font_size=12, color=ACCENT_RED, bold=True)

# ── Step 2 ──

add_section_header(slide, Inches(7.0), Inches(1.1), Inches(5),
                   "Step 2. 2차 항 (비대각 — 페널티)")

box = add_rounded_rect(slide, Inches(7.0), Inches(1.5), Inches(5.8), Inches(2.3),
                        SOFT_BLUE_BG, ACCENT_BLUE, Pt(1))
tb = add_text_box(slide, Inches(7.3), Inches(1.6), Inches(5.3), Inches(2.1))
tf = tb.text_frame
tf.word_wrap = True
set_text(tf, "모든 쌍 (i, j)에서, target 외 3가지 오답 각각에", font_size=14,
         color=DARK_GRAY)
add_paragraph(tf, "독립적인 양수 r을 뽑아 페널티 부여", font_size=14, color=DARK_GRAY,
              bold=True)
add_paragraph(tf, "", font_size=6)
add_paragraph(tf, "예) target pair (bᵢ, bⱼ) = (0, 0) 일 때:", font_size=13,
              color=ACCENT_BLUE, bold=True)
add_paragraph(tf, "", font_size=2)
add_paragraph(tf, " 오답(0,1): + r·(1−xᵢ)xⱼ    ← xᵢ=0,xⱼ=1일 때만 +r",
              font_size=13, color=DARK_GRAY, font_name="Consolas")
add_paragraph(tf, " 오답(1,0): + r·xᵢ(1−xⱼ)    ← xᵢ=1,xⱼ=0일 때만 +r",
              font_size=13, color=DARK_GRAY, font_name="Consolas")
add_paragraph(tf, " 오답(1,1): + r·xᵢxⱼ         ← xᵢ=1,xⱼ=1일 때만 +r",
              font_size=13, color=DARK_GRAY, font_name="Consolas")
add_paragraph(tf, "", font_size=4)
add_paragraph(tf, "target (0,0)에서 → 모든 페널티 = 0  ✓", font_size=13,
              color=ACCENT_GREEN, bold=True)

# ── 하단: 페널티 분배 테이블 ──

add_section_header(slide, Inches(0.6), Inches(4.1), Inches(6),
                   "페널티 함수의 Q 행렬 기여 (전개 결과)")

# 테이블
table_shape = slide.shapes.add_table(5, 5, Inches(0.6), Inches(4.5),
                                      Inches(9.5), Inches(2.1))
table = table_shape.table

# 열 너비
col_widths = [Inches(1.7), Inches(2.8), Inches(1.5), Inches(1.5), Inches(2.0)]
for i, w in enumerate(col_widths):
    table.columns[i].width = w

headers = ["오답 상태 (sᵢ, sⱼ)", "페널티 함수", "Qᵢᵢ 기여", "Qⱼⱼ 기여", "Qᵢⱼ 기여"]
rows_data = [
    ["(0, 0)", "r · (1−xᵢ)(1−xⱼ)", "−r", "−r", "+r"],
    ["(0, 1)", "r · (1−xᵢ)xⱼ", "0", "+r", "−r"],
    ["(1, 0)", "r · xᵢ(1−xⱼ)", "+r", "0", "−r"],
    ["(1, 1)", "r · xᵢxⱼ", "0", "0", "+r"],
]

# 헤더 스타일
for j, h in enumerate(headers):
    cell = table.cell(0, j)
    cell.text = h
    for p in cell.text_frame.paragraphs:
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.font.name = "맑은 고딕"
        p.alignment = PP_ALIGN.CENTER
    cell.fill.solid()
    cell.fill.fore_color.rgb = ACCENT_BLUE

# 데이터 스타일
for i, row in enumerate(rows_data):
    for j, val in enumerate(row):
        cell = table.cell(i + 1, j)
        cell.text = val
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(13)
            p.font.name = "Consolas"
            p.font.color.rgb = DARK_GRAY
            p.alignment = PP_ALIGN.CENTER
            if j >= 2 and val.startswith("+"):
                p.font.color.rgb = ACCENT_RED
            elif j >= 2 and val.startswith("−") or val.startswith("-"):
                p.font.color.rgb = ACCENT_BLUE
        cell.fill.solid()
        cell.fill.fore_color.rgb = WHITE if i % 2 == 0 else LIGHT_GRAY

# 우측 요약 박스
box = add_rounded_rect(slide, Inches(10.4), Inches(4.5), Inches(2.5), Inches(2.1),
                        SOFT_GREEN_BG, ACCENT_GREEN, Pt(1))
tb = add_text_box(slide, Inches(10.6), Inches(4.6), Inches(2.2), Inches(1.9))
tf = tb.text_frame
tf.word_wrap = True
set_text(tf, "보장", font_size=14, color=ACCENT_GREEN, bold=True,
         alignment=PP_ALIGN.CENTER)
add_paragraph(tf, "", font_size=4)
add_paragraph(tf, "target 상태에서", font_size=13, color=DARK_GRAY,
              alignment=PP_ALIGN.CENTER)
add_paragraph(tf, "모든 페널티 = 0", font_size=15, color=ACCENT_GREEN, bold=True,
              alignment=PP_ALIGN.CENTER)
add_paragraph(tf, "", font_size=4)
add_paragraph(tf, "오답 상태에서", font_size=13, color=DARK_GRAY,
              alignment=PP_ALIGN.CENTER)
add_paragraph(tf, "페널티 > 0", font_size=15, color=ACCENT_RED, bold=True,
              alignment=PP_ALIGN.CENTER)
add_paragraph(tf, "", font_size=4)
add_paragraph(tf, "∴ optimal = x*", font_size=16, color=BLACK, bold=True,
              alignment=PP_ALIGN.CENTER)

# 하단 의문
box = add_rounded_rect(slide, Inches(0.6), Inches(6.8), Inches(12.2), Inches(0.55),
                        SOFT_YELLOW_BG, ACCENT_ORANGE, Pt(1))
tb = add_text_box(slide, Inches(0.9), Inches(6.83), Inches(11.8), Inches(0.5))
set_text(tb.text_frame,
    "의문:  이렇게 만들어진 QUBO가 random하게 생성한 QUBO와 indistinguishable 할까?",
    font_size=15, color=ACCENT_ORANGE, bold=True, alignment=PP_ALIGN.CENTER)

add_slide_number(slide, 2, 3)


# ────────────────────────────────────────────────────────────────
# 슬라이드 3: 기댓값 0 달성
# ────────────────────────────────────────────────────────────────

slide = prs.slides.add_slide(prs.slide_layouts[6])
bg = slide.background.fill
bg.solid()
bg.fore_color.rgb = WHITE

add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), ACCENT_BLUE)

tb = add_text_box(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.6))
set_text(tb.text_frame, "3. 기댓값 0 달성 — 페널티 비율 최적화",
         font_size=26, color=BLACK, bold=True)

# ── 좌측: bᵢ=0일 때 분석 ──

add_section_header(slide, Inches(0.6), Inches(1.1), Inches(5),
                   "bᵢ = 0 일 때, xᵢ에 관련된 항 분석")

box = add_rounded_rect(slide, Inches(0.6), Inches(1.5), Inches(6.0), Inches(0.7),
                        LIGHT_GRAY)
tb = add_text_box(slide, Inches(0.9), Inches(1.55), Inches(5.5), Inches(0.6))
tf = tb.text_frame
tf.word_wrap = True
set_text(tf, "j ≠ i 인 변수들 중 bⱼ=0과 bⱼ=1이 각각 (n−1)/2 개씩 존재한다고 기대",
         font_size=13, color=DARK_GRAY)

# 분석 테이블
table_shape = slide.shapes.add_table(7, 5, Inches(0.6), Inches(2.4),
                                      Inches(6.0), Inches(3.0))
table = table_shape.table

col_w = [Inches(0.7), Inches(0.9), Inches(2.0), Inches(1.0), Inches(1.4)]
for i, w in enumerate(col_w):
    table.columns[i].width = w

t_headers = ["경우", "오답", "페널티 함수", "xᵢ 기여", "기대 횟수"]
t_rows = [
    ["bⱼ=0", "(0,1)", "r₀(xⱼ − xᵢxⱼ)", "0", "(n−1)/2"],
    ["bⱼ=0", "(1,0)", "r₁(xᵢ − xᵢxⱼ)", "+r₁", "(n−1)/2"],
    ["bⱼ=0", "(1,1)", "r₂(xᵢxⱼ)", "0", "(n−1)/2"],
    ["bⱼ=1", "(0,0)", "r₀'(1−xᵢ−xⱼ+xᵢxⱼ)", "−r₀'", "(n−1)/2"],
    ["bⱼ=1", "(1,0)", "r₁'(xᵢ − xᵢxⱼ)", "+r₁'", "(n−1)/2"],
    ["bⱼ=1", "(1,1)", "r₂'(xᵢxⱼ)", "0", "(n−1)/2"],
]

for j, h in enumerate(t_headers):
    cell = table.cell(0, j)
    cell.text = h
    for p in cell.text_frame.paragraphs:
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.font.name = "맑은 고딕"
        p.alignment = PP_ALIGN.CENTER
    cell.fill.solid()
    cell.fill.fore_color.rgb = ACCENT_BLUE

for i, row in enumerate(t_rows):
    for j, val in enumerate(row):
        cell = table.cell(i + 1, j)
        cell.text = val
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(12)
            p.font.name = "Consolas"
            p.font.color.rgb = DARK_GRAY
            p.alignment = PP_ALIGN.CENTER
            if j == 3 and val.startswith("+"):
                p.font.color.rgb = ACCENT_RED
                p.font.bold = True
            elif j == 3 and val.startswith("−"):
                p.font.color.rgb = ACCENT_BLUE
                p.font.bold = True
        cell.fill.solid()
        cell.fill.fore_color.rgb = WHITE if i % 2 == 0 else LIGHT_GRAY

# ── 우측: 기댓값 수식 + 해법 ──

add_section_header(slide, Inches(7.2), Inches(1.1), Inches(5), "기댓값 수식")

# xᵢ 계수 기댓값
box = add_rounded_rect(slide, Inches(7.2), Inches(1.5), Inches(5.6), Inches(1.5),
                        LIGHT_GRAY)
tb = add_text_box(slide, Inches(7.5), Inches(1.55), Inches(5.1), Inches(1.4))
tf = tb.text_frame
tf.word_wrap = True
set_text(tf, "기대 횟수가 동일하므로 (n−1)/2 로 약분:", font_size=12,
         color=MID_GRAY, italic=True)
add_paragraph(tf, "", font_size=4)
add_paragraph(tf, "xᵢ 계수 기댓값 ∝", font_size=13, color=ACCENT_BLUE, bold=True)
add_paragraph(tf,
    "  E[r₁] − E[r₀'] + E[r₁']",
    font_size=15, color=BLACK, bold=True, font_name="Consolas")
add_paragraph(tf, "", font_size=6)
add_paragraph(tf, "xᵢxⱼ 계수 기댓값 ∝", font_size=13, color=ACCENT_BLUE, bold=True)
add_paragraph(tf,
    "  −E[r₀] − E[r₁] + E[r₂]",
    font_size=14, color=BLACK, bold=True, font_name="Consolas")
add_paragraph(tf,
    "  +E[r₀'] − E[r₁'] + E[r₂']",
    font_size=14, color=BLACK, bold=True, font_name="Consolas")

# 해법 박스
add_section_header(slide, Inches(7.2), Inches(3.3), Inches(5), "해법", color=ACCENT_GREEN)

box = add_rounded_rect(slide, Inches(7.2), Inches(3.7), Inches(5.6), Inches(2.2),
                        SOFT_GREEN_BG, ACCENT_GREEN, Pt(1.5))
tb = add_text_box(slide, Inches(7.5), Inches(3.8), Inches(5.1), Inches(2.0))
tf = tb.text_frame
tf.word_wrap = True
set_text(tf, "E[r₀], E[r₀'] 를 다른 값들의 2배로 샘플링!", font_size=15,
         color=ACCENT_GREEN, bold=True)
add_paragraph(tf, "", font_size=4)
add_paragraph(tf, "E[r₁]=E[r₁']=E[r₂]=E[r₂']=k 로 놓으면:",
              font_size=12, color=MID_GRAY, italic=True)
add_paragraph(tf, "", font_size=2)
add_paragraph(tf, "  조건①: E[r₀'] = E[r₁]+E[r₁'] = 2k",
              font_size=13, color=DARK_GRAY, font_name="Consolas")
add_paragraph(tf, "  조건②: −E[r₀]+2k = 0  → E[r₀] = 2k",
              font_size=13, color=DARK_GRAY, font_name="Consolas")
add_paragraph(tf, "", font_size=4)
add_paragraph(tf, "예) 일반 r: Uniform(1,3) → E=2",
              font_size=13, color=DARK_GRAY, font_name="Consolas")
add_paragraph(tf, "   r₀,r₀': Uniform(2,6) → E=4 = 2×2  ✓",
              font_size=13, color=ACCENT_GREEN, bold=True, font_name="Consolas")

# 하단: 결론 & 한계
box = add_rounded_rect(slide, Inches(0.6), Inches(5.8), Inches(5.6), Inches(1.5),
                        SOFT_BLUE_BG, ACCENT_BLUE, Pt(1))
tb = add_text_box(slide, Inches(0.9), Inches(5.9), Inches(5.1), Inches(1.3))
tf = tb.text_frame
tf.word_wrap = True
set_text(tf, "달성한 것", font_size=14, color=ACCENT_BLUE, bold=True)
add_paragraph(tf, "", font_size=2)
add_paragraph(tf, "• E[qᵢⱼ] = 0: 비대각 항에서 target pair 정보 누출 차단",
              font_size=13, color=DARK_GRAY)
add_paragraph(tf, "• E[qᵢᵢ | bᵢ=0] = E[qᵢᵢ | bᵢ=1] = 0 도 동시에 달성",
              font_size=13, color=DARK_GRAY)
add_paragraph(tf, "• Q 행렬의 1차 통계량이 랜덤 QUBO와 동일",
              font_size=13, color=ACCENT_BLUE, bold=True)

box = add_rounded_rect(slide, Inches(7.2), Inches(6.15), Inches(5.6), Inches(1.15),
                        SOFT_RED_BG, ACCENT_RED, Pt(1))
tb = add_text_box(slide, Inches(7.5), Inches(6.2), Inches(5.1), Inches(1.0))
tf = tb.text_frame
tf.word_wrap = True
set_text(tf, "열린 문제", font_size=14, color=ACCENT_RED, bold=True)
add_paragraph(tf, "", font_size=2)
add_paragraph(tf, "• 1차 통계량(기댓값)은 위장 성공, 그러나 고차 통계량(분산, 상관)은?",
              font_size=13, color=DARK_GRAY)
add_paragraph(tf, "• SA-trivial: 에너지 지형 자체는 위장 못함 → 난이도 벤치마크로는 부족",
              font_size=13, color=DARK_GRAY)

add_slide_number(slide, 3, 3)


# ── 저장 ──

output_path = "/Users/yideun/Dev/qubo_dataset/zero_expectation/Zero_Expectation_방법론.pptx"
prs.save(output_path)
print(f"PPT 저장 완료: {output_path}")
